# import streamlit as st
# import pandas as pd
# import re
# import glob
# import io
# import os
# from pptx import Presentation
# from docx import Document
# from openai import AzureOpenAI


# # ============================================================
# # Helper Functions
# # ============================================================

# def call_llm(prompt, client, model_name):
#     """Call Azure OpenAI model."""
#     try:
#         response = client.chat.completions.create(
#             model=model_name,
#             messages=[{"role": "user", "content": prompt}],
#             temperature=0.4
#         )
#         return response.choices[0].message.content.strip()
#     except Exception as e:
#         return f"Error: {e}"


# def extract_ppt_text(ppt_path):
#     """Extract readable text from PPT."""
#     text = ""
#     prs = Presentation(ppt_path)
#     for slide in prs.slides:
#         for shape in slide.shapes:
#             if hasattr(shape, "text"):
#                 text += shape.text + "\n"
#     return text.strip()
# def extract_section_text(sow_text, section_number):
#             """
#             Extract text between numbered section headings like:
#             1. Executive Summary
#             2. Features of CoreAssess.AI
#             Works for headings with or without markdown (**, ##, etc.)
#             """
#             pattern = (
#                 rf"(?:\**\s*#*\s*)?"                # optional markdown prefix
#                 rf"{section_number}\.\s*[A-Za-z0-9 \-&()]+?(?:\**\s*)?"  # the numbered heading
#                 rf"(.*?)(?=\n\s*(?:\**\s*#*\s*)?\d+\.\s+[A-Z]|\Z)"      # stop at next section (digit + dot + capital)
#             )
#             match = re.search(pattern, sow_text, flags=re.DOTALL | re.IGNORECASE)
#             return match.group(1).strip() if match else ""

# # def extract_section_text(sow_text, section_number):
# #     """
# #     Extract text between sections like:
# #     **3. Key Findings & Recommendations Summary**
# #     Stops before the next numbered section (4., ## 4., etc.)
# #     """
# #     # Match the section header, like **3. Key Findings...**
# #     # and capture everything until the next numbered header.
# #     pattern = rf"{section_number}\.\s*[A-Za-z0-9 &–\-]+?(?=\s*\n)(.*?)(?=\n\s*(?:\**\s*#*\s*)?\d+\.\s|$)"

# #     match = re.search(pattern, sow_text, flags=re.DOTALL | re.IGNORECASE)
# #     if not match:
# #         # Try again with markdown-style headings like ## 3.
# #         pattern2 = rf"(?:\*\*|##)?\s*{section_number}\.\s*[A-Za-z0-9 &–\-]+?(?:\*\*|##)?\s*(.*?)(?=\n\s*(?:\*\*|##)?\s*\d+\.\s|$)"
# #         match = re.search(pattern2, sow_text, flags=re.DOTALL | re.IGNORECASE)

# #     return match.group(1).strip() if match else ""
# from docx.shared import Pt

# # def extract_section_text(sow_text, section_number):
# #     """
# #     Extract section text between numbered headings like:
# #     4. Benefits over Traditional Assessment
# #     Supports markdown (**, ##), tab spacing, and mixed symbols.
# #     """
# #     pattern = rf"(?is)(?:[#\*\s]*){section_number}\.\s*([A-Za-z0-9 \-&()]+?)[#\*\s:]*\n+(.*?)(?=\n\s*(?:[#\*\s]*)\d+\.\s|$)"
# #     match = re.search(pattern, sow_text)
# #     return match.group(2).strip() if match else ""
# # def extract_section_text(sow_text, section_number):
# #     """
# #     Extract section text between numbered headings like:
# #     2. Features of CoreAssess.AI
# #     Handles markdown (**), tabs, and spaces gracefully.
# #     """
# #     pattern = rf"(?is)(?:[#\*\s]*){section_number}\.\s*[A-Za-z0-9 \-&()]+?(?:\*\*|[#\s]*)?\s*(?:\r?\n)+(.*?)(?=\n\s*(?:[#\*\s]*)\d+\.\s|$)"
# #     match = re.search(pattern, sow_text)
# #     return match.group(1).strip() if match else ""



# def insert_annexure_table(doc, placeholder, df):
#     """Insert an Annexure-style table (Object, Issue, Key Modernization Steps) into the placeholder."""
#     inserted = False

#     for para in doc.paragraphs:
#         if placeholder in para.text:
#             inserted = True
#             para.text = ""  # clear placeholder text

#             # --- Create table ---
#             table = doc.add_table(rows=1, cols=3)
#             table.style = "Table Grid"

#             hdr_cells = table.rows[0].cells
#             hdr_cells[0].text = "Object Name"
#             hdr_cells[1].text = "Issue"
#             hdr_cells[2].text = "Key Modernization Steps"

#             # --- Format header cells ---
#             for cell in hdr_cells:
#                 for paragraph in cell.paragraphs:
#                     for run in paragraph.runs:
#                         run.bold = True
#                 cell.width = Pt(200)

#             # --- Populate rows from DataFrame ---
#             for _, row in df.iterrows():
#                 row_cells = table.add_row().cells
#                 row_cells[0].text = str(row.get("object name", row.get("Object Name", "")))
#                 row_cells[1].text = re.sub(r"<[^>]+>", "", str(row.get("issue", row.get("Issue", ""))))
#                 row_cells[2].text = re.sub(r"<[^>]+>", "", str(row.get("key modernization steps", row.get("Key Modernization Steps", ""))))

#             para._element.addnext(table._element)
#             break

#     if not inserted:
#         st.warning("⚠️ No <<ANNEXURE>> placeholder found. Appending Annexure at the end.")
#         doc.add_page_break()
#         doc.add_heading("Annexure — Modernization Object Summary", level=1)
#         table = doc.add_table(rows=1, cols=3)
#         table.style = "Table Grid"

#         hdr_cells = table.rows[0].cells
#         hdr_cells[0].text = "Object Name"
#         hdr_cells[1].text = "Issue"
#         hdr_cells[2].text = "Key Modernization Steps"

#         for _, row in df.iterrows():
#             row_cells = table.add_row().cells
#             row_cells[0].text = str(row.get("object name", row.get("Object Name", "")))
#             row_cells[1].text = re.sub(r"<[^>]+>", "", str(row.get("issue", row.get("Issue", ""))))
#             row_cells[2].text = re.sub(r"<[^>]+>", "", str(row.get("key modernization steps", row.get("Key Modernization Steps", ""))))

# # ============================================================
# # Core Function
# # ============================================================

# def generate_sow(df, client, model_name, client_name=None, repo_dir="Knowledge_Repo/Coreassess_KR"):
#     """Generate full SOW docx directly."""
#     client_ref = client_name if client_name else "the Client"

#     # Find available PPT references
#     ppt_files = glob.glob(os.path.join(repo_dir, "*.pptx"))
#     if not ppt_files:
#         ppt_text = "No PPTs found."
#         chosen_ppt = "None"
#     else:
#         ppt_text = extract_ppt_text(ppt_files[0])
#         chosen_ppt = ppt_files[0]

#     # Build prompt
#     # Build prompt
#     total = len(df)

#     # Safe extraction for sample issues
#     if "issue" in df.columns:
#         sample_col = df["issue"]
#     else:
#         sample_col = df.iloc[:, 0]

#     sample_issues = "; ".join(sample_col.astype(str).tolist()[:5])

#     prompt = f"""
#         You are a senior SAP consultant from Crave Infotech preparing a professional Statement of Work (SOW)
#         for a Clean Core Assessment (CoreAssess.AI) engagement with {client_ref}.

#         Below is Crave's official reference content from our CoreAssess knowledge presentation.
#         This content represents our internal tone, structure, and offering details.
#         Analyze it carefully to understand our standard messaging, flow, and technical vocabulary.

#         ---
#         {ppt_text}
#         ---

#         Now, using the reference as a guide (not to copy text directly), write a *comprehensive, polished, client-ready*
#         Statement of Work document that covers the following sections, aligned with Crave Infotech’s tone:

#         1. Executive Summary  
#         - Context of Clean Core Assessment  
#         - Value proposition of CoreAssess.AI  
#         - Alignment with SAP’s Clean Core strategy  


#         2. Features of CoreAssess.AI 
#         - Summarize the tool’s key technical capabilities and components.
#         - Mention capabilities like On-StackExtensibility, Side-by-Side Extensibility, SQL Analysis, and ROI Calculation.
#         - Explain each capability with 2-3 points
        
#         3. Key Findings & Recommendations Summary  
#         - Use aggregated counts from the data below  
#         - Present technical and business rationale for modernization

#         Total Objects: {total}  

#         Example Issues: {sample_issues}
#         - Include example insights and recommendations for each category — On-Stack, Side-by-Side, and Retire — drawn from the analyzed ABAP objects.  
#         - Use details from the data table to illustrate object-level examples, including Object Name, Issues, and Modernization Steps.  
#         - Structure like this:
#             1. On-Stack Extensibility: summarize key issues and modernization steps.
#             2. Side-by-Side Extensibility: summarize the findings and modernization actions.
#             3. Retire: summarize rationale and replacement steps.  
#         - End with a short summary paragraph connecting the recommendations to SAP’s Clean Core strategy.




#         5. Benefits over Traditional Assessment 
#         - Summarize the advantages of CoreAssess.AI compared to traditional clean core assessment methods.  
#         - Reflect the tone and structure from the reference PPT slides.  
#         - Highlight AI-dr

#         6. Working Together
#         - Use the tone and tiered structure from the “Working Together” slide in the PPT.  
#         - Describe Crave’s engagement options (Starter Pack, Silver, Gold, Platinum) in paragraph or tabular form.  
#         - Mention typical scope, duration, and pricing guidance (e.g., Complimentary, $50/object, $75/object, $95/object).  
#         - If the client provides a list of objects, mention that per-object pricing applies.  
#         - End with a short paragraph summarizing total estimated effort across 3 phases (Assessment, Recommendation, Presentation).

#         7. Working Together - ABAP Objects
#         - Use the reference tone from the following PPT section (do not copy directly):
#         - Summarize governance model, milestones, and deliverables.


#         **Important:**  
#         - Use Crave Infotech’s corporate tone — confident, concise, and consultative.  
#         - Do not reuse names or context from the reference PPT (like “Oatey Co.”).  
#         - Instead, personalize all context to {client_ref}.  
#         - Write complete paragraphs (not bullet slides).  
#         - Keep length around 4–6 pages of Word content.
#         """


#     # Get LLM result
#     # --- Split SOW by numbered headings like "1. Executive Summary" ---
#     full_sow = call_llm(prompt, client, model_name)





 

#     # Create a simple docx
#     # doc = Document()
#     # doc.add_heading("Statement of Work (SOW)", 0)
#     # doc.add_paragraph(full_sow.strip())

#     # # Save in memory
#     # buffer = io.BytesIO()
#     # doc.save(buffer)
#     # buffer.seek(0)

#     # st.success(f"✅ SOW generated using `{os.path.basename(chosen_ppt)}`")
#     # st.download_button(
#     #     label="📥 Download SOW Document (.docx)",
#     #     data=buffer,
#     #     file_name=f"SOW_{client_ref.replace(' ', '_')}.docx",
#     #     mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
#     # )
#     # --- Use Template ---
#     template_path = "Template/CoreAssess_Template.docx"

#     if os.path.exists(template_path):
#         doc = Document(template_path)
#         st.info("📄 Using Word template for SOW.")
#     else:
#         st.warning("⚠️ Template not found. Creating a blank document.")
#         doc = Document()

#     # --- Extract Executive Summary from SOW ---
#     # --- Extract Executive Summary (simple & forgiving) ---
#     # def extract_section_text(sow_text, section_number):
#     #     """
#     #     Extract text between numbered section headings like:
#     #     1. Executive Summary
#     #     2. Features of CoreAssess.AI
#     #     """
#     #     pattern = rf"(?:\**\s*#*\s*)?{section_number}\.\s*[A-Za-z0-9 \-&()]+?(?:\**\s*)?(.*?)(?=\n\s*(?:\**\s*#*\s*)?\d+\.\s|\Z)"
#     #     # pattern = rf"((?:\**\s*#*\s*)?{section_number}\.\s*[A-Za-z0-9 \-&()]+?(?:\**\s*)?(?:\n|$).*?)(?=\n\s*(?:\**\s*#*\s*)?\d+\.\s|\Z)"
#     #     # pattern = rf"(?s)(?:\*\*|\#\#)?\s*{section_number}\.\s*[A-Za-z0-9 \-&()]+?(?:\*\*|\#\#)?\s*(?:\r?\n)+.*?(?=(?:\r?\n)+\s*(?:\*\*|\#\#)?\s*\d+\.\s|\Z)"
#     #     match = re.search(pattern, sow_text, flags=re.DOTALL | re.IGNORECASE)
#     #     return match.group(1).strip() if match else ""


#         # match = re.search(pattern, sow_text, flags=re.DOTALL)
#         # if match:
#         #     # Extract both heading + body (group(0)) ensures title kept intact
#         #     section_text = match.group(0).strip()
#         #     return section_text
#         # return ""


#     # exec_match = re.search(
#     #     r"1\.\s*Executive\s*Summary(.*?)(?=\n\s*\d+\.\s|\Z)",
#     #     full_sow,
#     #     flags=re.DOTALL | re.IGNORECASE,
#     # )


#     # if exec_match:
#     #     exec_summary = exec_match.group(1).strip()
#     # else:
#     #     exec_summary = "Executive Summary not found in the generated SOW."
# # --- Extract sections we care about ---
#     exec_summary = extract_section_text(full_sow, 1)
#     features_section = extract_section_text(full_sow, 2)
#     findings_section = extract_section_text(full_sow, 3)
#     benefits_section = extract_section_text(full_sow, 4)
#     working_section = extract_section_text(full_sow, 5)
#     abap_section = extract_section_text(full_sow, 6)
#     insert_annexure_table(doc, "<<ANNEXURE>>", df)

#     if not exec_summary:
#         exec_summary = "Executive Summary not found in the generated SOW."
#     if not features_section:
#         features_section = "Features section not found in the generated SOW."

#     if not findings_section:
#         findings_section = "Findings section not found in the generated SOW."
#     if not benefits_section:
#         benefits_section = "Benefits section not found in the generated SOW."

#     if not working_section:
#         working_section = "Working Together section not found in the generated SOW."
#     if not abap_section:
#         abap_section = "ABAP Objects section not found in the generated SOW."
#     # --- Helper: insert text correctly (no upside-down order) ---
#     def insert_text_after_placeholder(doc, placeholder, text_block):
#         lines = [l.strip() for l in text_block.split("\n") if l.strip()]
#         inserted = False
#         for para in doc.paragraphs:
#             if placeholder in para.text:
#                 inserted = True
#                 para.text = ""  # clear placeholder text
#                 # reverse insertion to preserve order
#                 for line in reversed(lines):
#                     new_para = doc.add_paragraph(line)
#                     para._element.addnext(new_para._element)
#                 break
#         if not inserted:
#             st.warning(f"⚠️ Placeholder {placeholder} not found. Appending at end.")
#             doc.add_page_break()
#             doc.add_heading(placeholder.replace('<', '').replace('>', ''), level=1)
#             for line in lines:
#                 doc.add_paragraph(line)

#     # --- Insert sections ---
#     insert_text_after_placeholder(doc, "<<EXEC_SUMMARY>>", exec_summary)
#     insert_text_after_placeholder(doc, "<<FEATURES>>", features_section)
#     insert_text_after_placeholder(doc, "<<FINDINGS >>", findings_section)
#     insert_text_after_placeholder(doc, "<<BENEFITS>>", benefits_section)
#     insert_text_after_placeholder(doc, "<<WORKING TOGETHER>>", working_section)
#     insert_text_after_placeholder(doc, "<<ABAP OBJECTS>>", abap_section)
#     # --- Add full SOW at the end (so nothing is lost) ---
#     doc.add_page_break()
#     doc.add_heading("Full Statement of Work (Generated)", level=1)
#     for line in full_sow.split("\n"):
#         if line.strip():
#             doc.add_paragraph(line.rstrip())


#     # # --- Replace <<EXEC_SUMMARY>> in template ---
#     # placeholder_found = False
#     # for para in doc.paragraphs:
#     #     if "<<EXEC_SUMMARY>>" in para.text:
#     #         placeholder_found = True
#     #         # clear placeholder text
#     #         para.text = ""
#     #         # add the extracted Executive Summary content
#     #         # for line in exec_summary.split("\n"):
#     #         #     if line.strip():
#     #         #         new_para = doc.add_paragraph(line.strip())
#     #         #         para._element.addnext(new_para._element)
#     #         # Insert in correct order (so it doesn't appear upside down)
#     #         lines = [l.strip() for l in exec_summary.split("\n") if l.strip()]
#     #         for line in reversed(lines):  # reverse insertion order
#     #             new_para = doc.add_paragraph(line)
#     #             para._element.addnext(new_para._element)

#     #         break

#     # if not placeholder_found:
#     #     st.warning("⚠️ No <<EXEC_SUMMARY>> placeholder found. Appending at end of document.")
#     #     doc.add_paragraph(exec_summary)

#     # --- Add full SOW at the end ---
#     doc.add_page_break()
#     doc.add_heading("Full Statement of Work (Generated)", level=1)
#     doc.add_paragraph(full_sow.strip())

#     # --- Save in memory ---
#     buffer = io.BytesIO()
#     doc.save(buffer)
#     buffer.seek(0)

#     # --- Preview section ---
#     st.markdown("### 📄 Preview of Generated SOW")
#     preview_text = "\n".join(full_sow.split("\n")[:50])  # show first ~50 lines
#     st.text(preview_text.strip())


#     # --- Streamlit output ---
#     st.success(f"✅ SOW generated using `{os.path.basename(chosen_ppt)}` and inserted into template.")
#     st.download_button(
#         label="📥 Download SOW Document (.docx)",
#         data=buffer,
#         file_name=f"SOW_{client_ref.replace(' ', '_')}.docx",
#         mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
#     )



# # ============================================================
# # Streamlit UI
# # ============================================================

# def main():
#     st.title("🌐 CoreAssess.AI — Auto SOW Generator")

#     client_name = st.text_input("Client Name", placeholder="e.g., Adani Group")
#     uploaded = st.file_uploader("📂 Upload Excel (.xlsx)", type=["xlsx"])

#     if uploaded:
#         df = pd.read_excel(uploaded)
#         st.success(f"✅ File `{uploaded.name}` loaded successfully!")
#         st.dataframe(df.head(5))

#         # Azure OpenAI setup
#         client = AzureOpenAI(
#             azure_endpoint=os.getenv("AZURE_OPENAI_FRFP_ENDPOINT"),
#             api_key=os.getenv("AZURE_OPENAI_FRFP_KEY"),
#             api_version=os.getenv("AZURE_OPENAI_FRFP_VERSION")
#         )
#         model_name = "codetest"

#         if st.button("⚡ Generate SOW Document"):
#             generate_sow(df, client, model_name, client_name)


import streamlit as st
import pandas as pd
import re
import glob
import io
import os
from pptx import Presentation
from docx import Document
from openai import AzureOpenAI
from dotenv import load_dotenv


# --- Load your .env file safely ---
load_dotenv()

# # --- Normalize environment variables for Azure SDK ---
# # These three lines make sure AzureOpenAI gets what it expects
# os.environ["AZURE_OPENAI_API_KEY"] = os.getenv("AZURE_OPENAI_FRFP_KEY")
# os.environ["AZURE_OPENAI_ENDPOINT"] = "https://craveopenai.openai.azure.com/"
# os.environ["AZURE_OPENAI_API_VERSION"] = os.getenv("AZURE_OPENAI_FRFP_VERSION")

# # --- Initialize the Azure OpenAI client ---
# try:
#     client = AzureOpenAI(
#         azure_endpoint=os.environ["AZURE_OPENAI_ENDPOINT"],
#         api_key=os.environ["AZURE_OPENAI_API_KEY"],
#         api_version=os.environ["AZURE_OPENAI_API_VERSION"]
#     )
#     # st.info("✅ Connected to Azure OpenAI successfully.")
# except Exception as e:
#     st.error(f"⚠️ Azure OpenAI connection failed: {e}")


# ============================================================
# Helper Functions
# ============================================================

def call_llm(prompt, client, model_name):
    """Call Azure OpenAI model."""
    try:
        response = client.chat.completions.create(
            model=model_name,
            messages=[{"role": "user", "content": prompt}],
            temperature=0.4
        )
        return response.choices[0].message.content.strip()
    except Exception as e:
        return f"Error: {e}"



def extract_ppt_text(ppt_path):
    """
    Extract readable text from PPT (grouped shapes + tables) and detect both
    'Working Together' slides (Objects & ABAP Programs).
    """
    import re
    text = ""
    working_together_objects = ""
    working_together_abap = ""
    prs = Presentation(ppt_path)

    def extract_from_shape(shape):
        content = ""
        if hasattr(shape, "text") and shape.text.strip():
            content += shape.text.strip() + "\n"
        if hasattr(shape, "shapes"):  # recurse into grouped shapes
            for sub_shape in shape.shapes:
                content += extract_from_shape(sub_shape)
        if shape.shape_type == 19:  # handle tables
            for row in shape.table.rows:
                for cell in row.cells:
                    if cell.text.strip():
                        content += cell.text.strip() + "\n"
        return content

    for slide in prs.slides:
        slide_text = ""
        for shape in slide.shapes:
            slide_text += extract_from_shape(shape)
        clean_text = slide_text.strip()
        if not clean_text:
            continue

        if re.search(r"working\s*together", clean_text, re.IGNORECASE):
            if "objects" in clean_text.lower():
                working_together_objects = clean_text
            elif re.search(r"abap\s*program", clean_text, re.IGNORECASE):
                working_together_abap = clean_text

        text += clean_text + "\n\n"

    if working_together_objects or working_together_abap:
        st.success("✅ 'Working Together' slides successfully extracted from PPT.")
    else:
        st.warning("⚠️ Could not detect any 'Working Together' slides — check slide text formatting.")

    return text.strip(), working_together_objects.strip(), working_together_abap.strip()

from docx.shared import Pt

def insert_text(doc, heading_title, text_block):
    """Adds a page break, heading, and paragraphs from a text block."""
    if not text_block:
        return  # skip if empty
    doc.add_page_break()
    doc.add_heading(heading_title, level=1)
    for line in text_block.split("\n"):
        clean = line.strip()
        if not clean:
            continue
        p = doc.add_paragraph(clean)
        p_format = p.paragraph_format
        p_format.space_after = Pt(6)
        p_format.line_spacing = 1.2
        for run in p.runs:
            run.font.name = "Calibri"
            run.font.size = Pt(11)



def insert_annexure_table(doc, placeholder, df):
    """Insert an Annexure-style table (Object, Issue, Key Modernization Steps) into the placeholder."""
    inserted = False

    for para in doc.paragraphs:
        if placeholder in para.text:
            inserted = True
            para.text = ""  # clear placeholder text

            # --- Create table ---
            table = doc.add_table(rows=1, cols=3)
            table.style = "Table Grid"

            hdr_cells = table.rows[0].cells
            hdr_cells[0].text = "Object Name"
            hdr_cells[1].text = "Issue"
            hdr_cells[2].text = "Key Modernization Steps"

            # --- Format header cells ---
            for cell in hdr_cells:
                for paragraph in cell.paragraphs:
                    for run in paragraph.runs:
                        run.bold = True
                cell.width = Pt(200)

            # --- Populate rows from DataFrame ---
            for _, row in df.iterrows():
                row_cells = table.add_row().cells
                row_cells[0].text = str(row.get("object name", row.get("Object Name", "")))
                row_cells[1].text = re.sub(r"<[^>]+>", "", str(row.get("issue", row.get("Issue", ""))))
                row_cells[2].text = re.sub(r"<[^>]+>", "", str(row.get("key modernization steps", row.get("Key Modernization Steps", ""))))

            para._element.addnext(table._element)
            break

    if not inserted:
        st.warning("⚠️ No <<ANNEXURE>> placeholder found. Appending Annexure at the end.")
        doc.add_page_break()
        doc.add_heading("Annexure — Modernization Object Summary", level=1)
        table = doc.add_table(rows=1, cols=3)
        table.style = "Table Grid"

        hdr_cells = table.rows[0].cells
        hdr_cells[0].text = "Object Name"
        hdr_cells[1].text = "Issue"
        hdr_cells[2].text = "Key Modernization Steps"

        for _, row in df.iterrows():
            row_cells = table.add_row().cells
            row_cells[0].text = str(row.get("object name", row.get("Object Name", "")))
            row_cells[1].text = re.sub(r"<[^>]+>", "", str(row.get("issue", row.get("Issue", ""))))
            row_cells[2].text = re.sub(r"<[^>]+>", "", str(row.get("key modernization steps", row.get("Key Modernization Steps", ""))))

# ============================================================
# Core Function
# ============================================================

def generate_sow(df, client, model_name, client_name=None, repo_dir="Knowledge_Repo/Coreassess_KR"):
    """Generate full SOW docx directly."""
    client_ref = client_name if client_name else "the Client"

    # Find available PPT references
    ppt_files = glob.glob(os.path.join(repo_dir, "*.pptx"))
    if not ppt_files:
        ppt_text = "No PPTs found."
        chosen_ppt = "None"
    else:
        ppt_text, working_together_objects, working_together_abap = extract_ppt_text(ppt_files[0])

        chosen_ppt = ppt_files[0]


    # Build prompt
    total = len(df)

    # Safe extraction for sample issues
    if "issue" in df.columns:
        sample_col = df["issue"]
    else:
        sample_col = df.iloc[:, 0]

    sample_issues = "; ".join(sample_col.astype(str).tolist()[:5])

    prompt = f"""

        IMPORTANT: Do NOT include any top-level title, cover line, or "Statement of Work" header.
        Start the document with "1. Executive Summary" as the very first line.

        You are a senior SAP consultant from Crave Infotech preparing a professional Statement of Work (SOW)
        for a Clean Core Assessment (CoreAssess.AI) engagement with {client_ref}.

        Below is Crave's official reference content from our CoreAssess knowledge presentation.
        This content represents our internal tone, structure, and offering details.
        Analyze it carefully to understand our standard messaging, flow, and technical vocabulary.

        ---
        {ppt_text}
        ---

        Now, using the reference as a guide (not to copy text directly), write a *comprehensive, polished, client-ready*
        Statement of Work document that covers the following sections, aligned with Crave Infotech’s tone:

        1. Executive Summary  
        - Context of Clean Core Assessment  
        - Value proposition of CoreAssess.AI  
        - Alignment with SAP’s Clean Core strategy  


        2. Features of CoreAssess.AI 
        - Summarize the tool’s key technical capabilities and components.
        - Mention capabilities like On-StackExtensibility, Side-by-Side Extensibility, SQL Analysis, and ROI Calculation.
        - Explain each capability with 2-3 points
        
        3. Key Insights & Recommendations  
        - Use aggregated counts from the data below  
        - Present technical and business rationale for modernization

        Total Objects: {total}  

        Example Issues: {sample_issues}
        - Include example insights and recommendations for each category — On-Stack, Side-by-Side, and Retire — drawn from the analyzed ABAP objects.  
        - Use details from the data table to illustrate object-level examples, including Object Name, Issues, and Modernization Steps.  
        - Structure like this:
            1. On-Stack Extensibility: summarize key issues and modernization steps.
            2. Side-by-Side Extensibility: summarize the findings and modernization actions.
            3. Retire: summarize rationale and replacement steps.  
        - End with a short summary paragraph connecting the recommendations to SAP’s Clean Core strategy.




        5. Benefits over Traditional Assessment 
        - Summarize the advantages of CoreAssess.AI compared to traditional clean core assessment methods.  
        - Reflect the tone and structure from the reference PPT slides.  
        - Highlight AI-dr

        6. Working Together
        - Use the tone and tiered structure from the “Working Together” slide in the PPT.  
        - Describe Crave’s engagement options (Starter Pack, Silver, Gold, Platinum) in paragraph or tabular form.  
        - Mention typical scope, duration, and pricing guidance (e.g., Complimentary, $10/object, $7.5/object, $5/object).
        - If the client provides a list of objects, mention that per-object pricing applies.  
        - End with a short paragraph summarizing total estimated effort across 3 phases (Assessment, Recommendation, Presentation).

        7. Working Together - ABAP Objects
        - Use the tone and tiered structure from the “Working Together” slide in the PPT.  
        - Describe Crave’s engagement options (Starter Pack, Silver, Gold, Platinum) in paragraph or tabular form.  
        - Mention typical scope, duration, and pricing guidance (e.g., Complimentary, $100/Program, $75/Program, $50/Program).
        - If the client provides a list of objects, mention that per-object pricing applies.  
        - End with a short paragraph summarizing total estimated effort across 3 phases (Assessment, Recommendation, Presentation).


        **Important:**  
        - Use Crave Infotech’s corporate tone — confident, concise, and consultative.  
        - Do not reuse names or context from the reference PPT (like “Oatey Co.”).  
        - Instead, personalize all context to {client_ref}.  
        - Write complete paragraphs (not bullet slides).  
        - Keep length around 4–6 pages of Word content.
        """


    # Get LLM result
    # --- Split SOW by numbered headings like "1. Executive Summary" ---
    full_sow = call_llm(prompt, client, model_name)


    # --- Use Template ---
    template_path = "Template/CoreAssess_Template.docx"

    if os.path.exists(template_path):
        doc = Document(template_path)
        st.info("📄 Using Word template for SOW.")
    else:
        st.warning("⚠️ Template not found. Creating a blank document.")
        doc = Document()

    # def insert_full_sow(doc, placeholder, sow_text):
    #     """Insert AI-generated text into template and add page breaks for top-level sections as Heading 1."""
    #     inserted = False
    #     for para in doc.paragraphs:
    #         if placeholder in para.text:
    #             inserted = True
    #             para.text = ""  # clear placeholder text

    #             lines = [line.strip() for line in sow_text.split("\n") if line.strip()]
    #             for line in lines:
    #                 # ✅ Top-level section (e.g. "1. Executive Summary") → page break + Heading 1
    #                 if re.match(r"^\s*\d+\.\s+[A-Z]", line):
    #                     doc.add_page_break()
    #                     # doc.add_heading(line, level=1)
    #                     heading_text = re.sub(r"^\s*\d+\.\s*", "", line).strip()
    #                     doc.add_heading(heading_text, level=1)
                    

    #                 else:
    #                     p = doc.add_paragraph(line)
    #                     # Standard formatting for body paragraphs
    #                     p_format = p.paragraph_format
    #                     p_format.space_after = Pt(6)
    #                     p_format.line_spacing = 1.2
    #                     for run in p.runs:
    #                         run.font.name = "Calibri"
    #                         run.font.size = Pt(11)
    #             break
            
    #     if not inserted:
    #         st.warning("⚠️ Placeholder <<CONTENT START>> not found. Appending at end.")
    #         doc.add_page_break()
    #         for line in sow_text.split("\n"):
    #             if line.strip():
    #                 doc.add_paragraph(line.strip())
    def insert_full_sow(doc, placeholder, sow_text):
        """Insert AI-generated SOW content and format sections with headings & bullets."""
        inserted = False
        for para in doc.paragraphs:
            if placeholder in para.text:
                inserted = True
                para.text = ""  # clear placeholder text

                lines = [line.strip() for line in sow_text.split("\n") if line.strip()]
                last_heading = None

                for line in lines:

                    # === Heading 1 (Top-level sections like "1. Executive Summary") ===
                    if re.match(r"^\s*\d+\.\s+[A-Z]", line):
                        heading_text = re.sub(r"^\s*\d+\.\s*", "", line).strip()
                        if heading_text != last_heading:  # avoid repeats
                            doc.add_page_break()
                            doc.add_heading(heading_text, level=1)
                            last_heading = heading_text
                        continue  # ✅ stop further processing of this line

                    # === Bullet sub-heading like "- **On-Stack Extensibility**:" ===
                    # if re.match(r"^-\s*\*\*[A-Za-z \-]+\*\*:", line):
                    #     heading_text = re.sub(r"^-\s*", "", line).strip()
                    #     try:
                    #         doc.add_paragraph(heading_text, style="Crave Heading_4")
                    #     except KeyError:
                    #         # fallback if custom style not found
                    #         p = doc.add_paragraph()
                    #         run = p.add_run(heading_text)
                    #         run.bold = True
                    #         run.font.name = "Calibri"
                    #         run.font.size = Pt(11)
                    #     continue
                    # === Bullet sub-heading like "- **On-Stack Extensibility**:" ===
                    # if re.match(r"^-\s*\*\*[A-Za-z \-]+\*\*:", line):
                    #     # remove bullet (-), bold markdown (**), and colon (:)
                    #     heading_text = re.sub(r"^-\s*|\*\*|:", "", line).strip()
                    #     try:
                    #         doc.add_paragraph(heading_text, style="Crave Heading_4")
                    #     except KeyError:
                    #         # fallback if custom style not found
                    #         p = doc.add_paragraph()
                    #         run = p.add_run(heading_text)
                    #         run.bold = True
                    #         run.font.name = "Calibri"
                    #         run.font.size = Pt(11)
                    #     continue
                    # === True section subheadings (like On-Stack, Side-by-Side, SQL Analysis, ROI Calculation) ===
                    # === True subheadings: technical or engagement tiers ===
                    if re.match(
                        # r"^(-|\d+\.)\s*\*\*(On-Stack|Side-by-Side|SQL|ROI|Starter Pack|Silver|Gold|Platinum)[A-Za-z \-]*\*\*:",
                        r"^(-|\d+\.)\s*\*\*(On-Stack|Side-by-Side|SQL|ROI|Starter Pack|Silver|Gold|Platinum|Retire)[A-Za-z \-]*\*\*:",
                        line,
                    ):
                        heading_text = re.sub(r"^(-|\d+\.)\s*|\*\*|:", "", line).strip()
                        try:
                            doc.add_paragraph(heading_text, style="Crave Heading_4")
                        except KeyError:
                            p = doc.add_paragraph()
                            run = p.add_run(heading_text)
                            run.bold = True
                            run.font.name = "Calibri"
                            run.font.size = Pt(11)
                        continue



                    # === Inline label bolding (like "**Object Name**: ZOTCR_DEL_MULTI_TRAILERS") ===
                    if re.match(r"^-\s*\*\*[A-Za-z ]+\*\*:", line):
                        # Extract label (inside ** **) and the rest of the content
                        match = re.match(r"^-\s*\*\*([A-Za-z ]+)\*\*:\s*(.*)", line)
                        if match:
                            label, content = match.groups()
                            p = doc.add_paragraph()
                            run_label = p.add_run(f"{label}: ")
                            run_label.bold = True
                            run_label.font.name = "Calibri"
                            run_label.font.size = Pt(11)

                            if content:
                                run_content = p.add_run(content.strip())
                                run_content.font.name = "Calibri"
                                run_content.font.size = Pt(11)
                            continue


                        # === Normal bullet points (e.g. "- Provides high-level recommendations...") ===
                        # === Simple markdown-style bullets (works reliably with Word templates) ===
                    if line.startswith("- ") or line.startswith("• "):
                            # remove bullet marker
                            bullet_text = line[2:].strip() if line.startswith("- ") else line[1:].strip()

                            # add as List Bullet 2 (or fallback)
                            try:
                                p = doc.add_paragraph(bullet_text, style="List Bullet 2")
                            except KeyError:
                                # fallback if the style is missing in template
                                p = doc.add_paragraph(f"• {bullet_text}")

                            # formatting
                            p.paragraph_format.left_indent = Pt(18)
                            p.paragraph_format.space_after = Pt(2)
                            p.paragraph_format.line_spacing = 1.15

                            for run in p.runs:
                                run.font.name = "Calibri"
                                run.font.size = Pt(11)
                            continue


                    # === Regular paragraph text ===
                    p = doc.add_paragraph(line)
                    p_format = p.paragraph_format
                    p_format.space_after = Pt(6)
                    p_format.line_spacing = 1.2
                    for run in p.runs:
                        run.font.name = "Calibri"
                        run.font.size = Pt(11)

                break  # finished inserting

        # === If placeholder not found ===
        if not inserted:
            st.warning("⚠️ Placeholder <<CONTENT START>> not found. Appending at end.")
            doc.add_page_break()
            for line in sow_text.split("\n"):
                if line.strip():
                    doc.add_paragraph(line.strip())



    # --- Insert generated content ---
    insert_full_sow(doc, "<<CONTENT START>>", full_sow)


    # --- Add Annexure section at the end ---
    doc.add_page_break()
    doc.add_heading("Annexure — Modernization Object Summary", level=1)

    # Create table with 3 columns: Object, Issue, Modernization Steps
    table = doc.add_table(rows=1, cols=3)
    table.style = "Table Grid"

    hdr_cells = table.rows[0].cells
    hdr_cells[0].text = "Object Name"
    hdr_cells[1].text = "Issue"
    hdr_cells[2].text = "Key Modernization Steps"

    # Format header
    for cell in hdr_cells:
        for paragraph in cell.paragraphs:
            for run in paragraph.runs:
                run.bold = True

    # Fill data from uploaded Excel (df)
    for _, row in df.iterrows():
        row_cells = table.add_row().cells
        row_cells[0].text = str(row.get("object name", row.get("Object Name", "")))
        row_cells[1].text = str(row.get("issue", row.get("Issue", "")))
        row_cells[2].text = str(row.get("key modernization steps", row.get("Key Modernization Steps", "")))

    # Add a closing note
    doc.add_paragraph()
    doc.add_paragraph("This annexure provides a detailed mapping of identified objects, their issues, and the corresponding modernization steps proposed by CoreAssess.AI.")


    # --- Save to memory ---
    buffer = io.BytesIO()
    doc.save(buffer)
    buffer.seek(0)

    # --- Preview + Download ---
    st.markdown("### 📄 Preview of Generated SOW")
    preview_text = "\n".join(full_sow.split("\n")[:50])
    st.text(preview_text.strip())

    st.success(f"✅ SOW generated using `{os.path.basename(chosen_ppt)}` and inserted into template.")
    st.download_button(
        label="📥 Download SOW Document (.docx)",
        data=buffer,
        file_name=f"SOW_{client_ref.replace(' ', '_')}.docx",
        mime="application/vnd.openxmlformats-officedocument.wordprocessingml.document"
    )

# ============================================================
# Streamlit UI
# ============================================================

def main():
    st.title("🌐 CoreAssess.AI — Auto SOW Generator")

    client_name = st.text_input("Client Name", placeholder="e.g., Adani Group")
    uploaded = st.file_uploader("📂 Upload Excel (.xlsx)", type=["xlsx"])

    if uploaded:
        df = pd.read_excel(uploaded)
        st.success(f"✅ File `{uploaded.name}` loaded successfully!")
        st.dataframe(df.head(5))

        # Azure OpenAI setup
        client = AzureOpenAI(
            azure_endpoint=os.getenv("AZURE_OPENAI_FRFP_ENDPOINT"),
            api_key=os.getenv("AZURE_OPENAI_FRFP_KEY"),
            api_version=os.getenv("AZURE_OPENAI_FRFP_VERSION")
        )
        model_name = "codetest"

        if st.button("⚡ Generate SOW Document"):
            generate_sow(df, client, model_name, client_name)